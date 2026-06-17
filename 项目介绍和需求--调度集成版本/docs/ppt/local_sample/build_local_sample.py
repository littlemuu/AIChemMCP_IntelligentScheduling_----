from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_PATH = OUT_DIR / "sample_preview.pptx"


NAVY = RGBColor(11, 31, 58)
NAVY_2 = RGBColor(17, 46, 83)
BLUE = RGBColor(30, 107, 255)
CYAN = RGBColor(22, 199, 217)
DARK = RGBColor(43, 52, 64)
MID = RGBColor(92, 108, 128)
LIGHT_BG = RGBColor(244, 248, 252)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(228, 240, 249)
GREEN = RGBColor(42, 181, 119)
ORANGE = RGBColor(245, 146, 64)


def inch(value):
    return Inches(value)


def set_font(run, size, color=DARK, bold=False, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.05)
    tf.margin_right = inch(0.05)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return box


def add_box(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.1)
    return shape


def add_labeled_box(slide, x, y, w, h, title, subtitle="", fill=WHITE, line=PALE,
                    title_color=DARK, subtitle_color=MID, title_size=15):
    add_box(slide, x, y, w, h, fill, line)
    add_textbox(slide, x + 0.14, y + 0.13, w - 0.28, 0.26, title, title_size, title_color, True)
    if subtitle:
        add_textbox(slide, x + 0.14, y + 0.46, w - 0.28, h - 0.55, subtitle, 9.5, subtitle_color)


def add_chip(slide, x, y, text, fill=CYAN, color=WHITE, w=None):
    width = w if w is not None else 0.22 + len(text) * 0.12
    add_box(slide, x, y, width, 0.34, fill, None)
    add_textbox(slide, x + 0.08, y + 0.065, width - 0.16, 0.18, text, 8.5, color, True, PP_ALIGN.CENTER)


def add_arrow(slide, x, y, w, h=0.25, color=CYAN):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(x), inch(y), inch(w), inch(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_small_label(slide, x, y, text, color=MID):
    add_textbox(slide, x, y, 2.2, 0.22, text, 8, color, False)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="AIChemMCP scheduling prototype | simulation / mock hardware"):
    add_textbox(slide, 0.55, 7.08, 7.2, 0.22, text, 7.5, MID)


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Layered local "tech" background, all editable shapes.
    add_box(slide, 0, 0, 13.33, 0.18, CYAN, None, False)
    add_box(slide, 0, 6.72, 13.33, 0.78, NAVY_2, None, False)
    for i in range(10):
        x = 7.2 + i * 0.58
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(0.55), inch(0.015), inch(5.8))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(24, 78, 120)
        line.line.fill.background()
    for i in range(6):
        y = 0.8 + i * 0.78
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(7.0), inch(y), inch(5.9), inch(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(24, 78, 120)
        line.line.fill.background()

    add_textbox(slide, 0.68, 0.72, 4.8, 0.34, "工程实践结题汇报", 13, CYAN, True)
    add_textbox(slide, 0.62, 1.35, 6.4, 1.1, "AIChemMCP\n智能调度原型系统", 34, WHITE, True)
    add_textbox(slide, 0.68, 2.78, 5.55, 0.42, "面向自动化化学实验流程的仿真调度与模拟设备接口原型", 14, RGBColor(210, 226, 240))
    add_chip(slide, 0.68, 3.58, "Action Server", BLUE, WHITE, 1.45)
    add_chip(slide, 2.28, 3.58, "Scheduling Runtime", CYAN, NAVY, 1.72)
    add_chip(slide, 4.18, 3.58, "Simulation / Mock", RGBColor(44, 92, 133), WHITE, 1.58)

    # Abstract system/lab motif.
    add_labeled_box(slide, 7.55, 1.12, 2.05, 0.86, "Agent", "工具调用入口", RGBColor(21, 57, 96), CYAN, WHITE, RGBColor(190, 220, 235), 15)
    add_labeled_box(slide, 10.25, 1.12, 2.12, 0.86, "Action Server", "6 个工具接口", RGBColor(21, 57, 96), CYAN, WHITE, RGBColor(190, 220, 235), 15)
    add_arrow(slide, 9.58, 1.42, 0.58, 0.18, CYAN)
    add_labeled_box(slide, 8.2, 2.68, 3.6, 0.96, "SchedulingRuntime", "任务队列 · 状态快照 · 执行历史", RGBColor(14, 72, 116), CYAN, WHITE, RGBColor(190, 220, 235), 15)
    add_labeled_box(slide, 7.15, 4.35, 1.7, 0.72, "WS_REACTOR_A", "reaction", RGBColor(28, 84, 128), None, WHITE, RGBColor(190, 220, 235), 10)
    add_labeled_box(slide, 9.15, 4.35, 1.7, 0.72, "WS_MEASURE_A", "yield / pH", RGBColor(28, 84, 128), None, WHITE, RGBColor(190, 220, 235), 10)
    add_labeled_box(slide, 11.15, 4.35, 1.7, 0.72, "WS_CHAR_A", "HPLC / NMR", RGBColor(28, 84, 128), None, WHITE, RGBColor(190, 220, 235), 10)
    add_textbox(slide, 7.4, 5.62, 4.8, 0.28, "当前定位：仿真调度原型系统，未接入真实实验室硬件", 10.5, RGBColor(210, 226, 240), False, PP_ALIGN.CENTER)

    add_textbox(slide, 0.68, 6.86, 4.4, 0.25, "素材依据：README.md / docs/ppt_materials.md / demo_result.json", 8, RGBColor(187, 205, 220))
    return slide


def make_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_box(slide, 0, 0, 13.33, 0.18, CYAN, None, False)
    add_textbox(slide, 0.55, 0.42, 3.8, 0.42, "系统总体架构", 24, NAVY, True)
    add_textbox(slide, 8.4, 0.48, 4.2, 0.22, "接口层与调度层解耦，形成可运行的仿真闭环", 10, MID, False, PP_ALIGN.RIGHT)

    y = 1.55
    boxes = [
        (0.55, "用户 / Agent", "agent.py\n本地演示入口"),
        (2.62, "Action Server", "servers/action_server.py\n工具声明与路由"),
        (4.92, "Tools 工具封装", "tools/action_server_tools.py\n参数转换与错误包装"),
        (7.25, "SchedulingRuntime", "scheduling/runtime.py\n任务 · 时间 · 状态 · 历史"),
        (9.75, "Scheduler", "scheduling/scheduler.py\n资源时间线预留"),
        (11.55, "Models", "Task / Resource\nWorkstation / Robot"),
    ]
    widths = [1.55, 1.85, 1.95, 2.0, 1.45, 1.35]
    for idx, (x, title, subtitle) in enumerate(boxes):
        fill = WHITE if idx < 3 else RGBColor(232, 246, 250)
        line = BLUE if idx in (1, 3, 4) else RGBColor(190, 210, 226)
        add_labeled_box(slide, x, y, widths[idx], 1.22, title, subtitle, fill, line, NAVY, MID, 12)
        if idx < len(boxes) - 1:
            add_arrow(slide, x + widths[idx] + 0.08, y + 0.49, 0.35, 0.18, CYAN)

    # Surrounding inputs and outputs.
    add_labeled_box(slide, 0.88, 3.65, 2.15, 0.82, "examples/*.json", "示例资源、任务、三步 workflow", WHITE, RGBColor(190, 210, 226), NAVY, MID, 12)
    add_arrow(slide, 2.55, 3.22, 0.9, 0.18, RGBColor(109, 180, 210))
    add_labeled_box(slide, 5.32, 3.58, 2.65, 0.95, "Action Tools", "robotic_reaction / measurement\ncharacterization / status / advance", WHITE, BLUE, NAVY, MID, 12)
    add_labeled_box(slide, 9.72, 3.65, 2.55, 0.82, "outputs/demo_result.json", "结构化 demo 输出与资源使用记录", WHITE, RGBColor(190, 210, 226), NAVY, MID, 12)
    add_arrow(slide, 8.1, 3.98, 1.3, 0.18, RGBColor(109, 180, 210))

    add_labeled_box(slide, 0.76, 5.25, 2.9, 0.82, "接口层", "暴露 6 个可调用工具，屏蔽调度内部细节", RGBColor(229, 238, 252), BLUE, NAVY, MID, 13)
    add_labeled_box(slide, 4.02, 5.25, 3.0, 0.82, "运行时层", "维护任务队列、仿真时间、状态快照、执行历史", RGBColor(226, 248, 250), CYAN, NAVY, MID, 13)
    add_labeled_box(slide, 7.38, 5.25, 2.9, 0.82, "调度层", "工作站 / 机器人时间线预留与状态推进", RGBColor(226, 248, 250), CYAN, NAVY, MID, 13)
    add_labeled_box(slide, 10.62, 5.25, 1.95, 0.82, "验证", "tests + demo smoke", WHITE, GREEN, NAVY, MID, 13)

    add_footer(slide)
    return slide


def make_demo_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_box(slide, 0, 0, 13.33, 0.18, CYAN, None, False)
    add_textbox(slide, 0.55, 0.42, 5.5, 0.42, "Demo 演示流程 / 调度闭环", 24, NAVY, True)
    add_textbox(slide, 8.3, 0.48, 4.35, 0.22, "真实输出依据：AIChemMCP-main/outputs/demo_result.json", 9.2, MID, False, PP_ALIGN.RIGHT)

    # Timeline.
    y = 1.45
    stages = [
        ("提交任务", "time 0-4\n4 个任务进入运行时"),
        ("推进时间", "time 39\nTASK-0001 完成"),
        ("运行到完成", "time 152\nall_completed: true"),
    ]
    x_positions = [0.78, 4.95, 9.12]
    for i, (title, subtitle) in enumerate(stages):
        add_labeled_box(slide, x_positions[i], y, 2.75, 0.98, title, subtitle, WHITE, BLUE if i == 0 else CYAN, NAVY, MID, 15)
        if i < 2:
            add_arrow(slide, x_positions[i] + 2.9, y + 0.39, 1.1, 0.18, CYAN)

    # Closed loop visual.
    add_labeled_box(slide, 0.78, 3.15, 2.25, 0.72, "工具调用", "robotic_* / scheduler_*", RGBColor(232, 246, 250), CYAN, NAVY, MID, 12)
    add_labeled_box(slide, 3.55, 3.15, 2.25, 0.72, "Runtime", "任务、时间、状态", RGBColor(232, 246, 250), CYAN, NAVY, MID, 12)
    add_labeled_box(slide, 6.32, 3.15, 2.25, 0.72, "Scheduler", "资源时间线预留", RGBColor(232, 246, 250), CYAN, NAVY, MID, 12)
    add_labeled_box(slide, 9.09, 3.15, 2.25, 0.72, "状态快照", "runtime_status", RGBColor(232, 246, 250), CYAN, NAVY, MID, 12)
    add_labeled_box(slide, 5.02, 4.38, 2.8, 0.72, "demo_result.json", "结构化输出归档", WHITE, BLUE, NAVY, MID, 12)
    add_arrow(slide, 3.05, 3.39, 0.42, 0.16, CYAN)
    add_arrow(slide, 5.82, 3.39, 0.42, 0.16, CYAN)
    add_arrow(slide, 8.59, 3.39, 0.42, 0.16, CYAN)
    add_arrow(slide, 6.12, 4.02, 0.78, 0.15, RGBColor(109, 180, 210))

    # Evidence cards.
    add_labeled_box(slide, 0.78, 5.72, 2.1, 0.72, "4 tasks", "TASK-0001 ~ TASK-0004", WHITE, RGBColor(190, 210, 226), NAVY, MID, 14)
    add_labeled_box(slide, 3.15, 5.72, 2.1, 0.72, "3-step workflow", "reaction -> yield -> hplc", WHITE, RGBColor(190, 210, 226), NAVY, MID, 14)
    add_labeled_box(slide, 5.52, 5.72, 2.1, 0.72, "simulation_time", "152", WHITE, RGBColor(190, 210, 226), NAVY, MID, 14)
    add_labeled_box(slide, 7.89, 5.72, 2.1, 0.72, "all_completed", "true", RGBColor(230, 250, 240), GREEN, NAVY, GREEN, 14)
    add_labeled_box(slide, 10.26, 5.72, 2.1, 0.72, "tests", "8 unittest OK", WHITE, RGBColor(190, 210, 226), NAVY, MID, 14)

    add_textbox(slide, 0.78, 6.68, 5.8, 0.24, "说明：当前演示为 simulation / mock hardware，不声明真实硬件接入。", 8.5, MID)
    add_footer(slide)
    return slide


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = inch(13.333333)
    prs.slide_height = inch(7.5)
    make_cover(prs)
    make_architecture(prs)
    make_demo_loop(prs)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build()
