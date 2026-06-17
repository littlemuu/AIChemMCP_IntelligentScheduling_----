from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_PATH = OUT_DIR / "sample_preview_v2.pptx"


NAVY = RGBColor(9, 26, 48)
NAVY_2 = RGBColor(12, 45, 78)
BLUE = RGBColor(30, 107, 255)
CYAN = RGBColor(18, 196, 218)
DARK = RGBColor(37, 49, 64)
MID = RGBColor(91, 108, 128)
LIGHT_BG = RGBColor(246, 249, 253)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(226, 240, 249)
PALE_CYAN = RGBColor(229, 250, 252)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(245, 146, 64)


def inch(value):
    return Inches(value)


def set_font(run, size, color=DARK, bold=False, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return box


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, width=1.0):
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    return shape


def add_round(slide, x, y, w, h, fill=WHITE, line=PALE, width=1.0):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, width)


def add_rect(slide, x, y, w, h, fill, line=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)


def add_title(slide, title, subtitle=None, dark=False):
    title_color = WHITE if dark else NAVY
    sub_color = RGBColor(190, 212, 230) if dark else MID
    add_text(slide, 0.62, 0.38, 5.2, 0.44, title, 26, title_color, True)
    if subtitle:
        add_text(slide, 0.66, 0.88, 7.2, 0.26, subtitle, 10.5, sub_color)


def add_footer(slide, dark=False):
    color = RGBColor(183, 207, 224) if dark else MID
    add_text(slide, 0.62, 7.12, 6.8, 0.22,
             "AIChemMCP scheduling prototype | simulation / mock hardware", 7.5, color)


def add_chip(slide, x, y, w, text, fill, color=WHITE):
    add_round(slide, x, y, w, 0.34, fill, None)
    add_text(slide, x + 0.07, y + 0.07, w - 0.14, 0.16, text, 8.5, color, True, PP_ALIGN.CENTER)


def add_main_card(slide, x, y, w, h, title, role, path, accent=CYAN):
    add_round(slide, x + 0.04, y + 0.04, w, h, RGBColor(215, 225, 236), None)
    add_round(slide, x, y, w, h, WHITE, RGBColor(191, 214, 230), 1.1)
    add_rect(slide, x, y, 0.12, h, accent, None)
    add_text(slide, x + 0.22, y + 0.15, w - 0.34, 0.22, title, 11.5, NAVY, True)
    add_text(slide, x + 0.22, y + 0.43, w - 0.34, 0.20, role, 8.8, DARK)
    add_text(slide, x + 0.22, y + 0.68, w - 0.34, 0.18, path, 7.0, MID)


def add_aux_card(slide, x, y, w, h, title, line1, line2, accent=BLUE):
    add_round(slide, x, y, w, h, WHITE, RGBColor(202, 221, 235), 0.9)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.16, y + 0.17, 0.25, 0.25, accent, None)
    add_text(slide, x + 0.50, y + 0.13, w - 0.64, 0.22, title, 11, NAVY, True)
    add_text(slide, x + 0.50, y + 0.43, w - 0.64, 0.18, line1, 8.8, DARK)
    add_text(slide, x + 0.50, y + 0.67, w - 0.64, 0.18, line2, 8.2, MID)


def add_arrow(slide, x, y, w, h=0.22, color=CYAN):
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h, color, None)
    return arrow


def add_line(slide, x1, y1, x2, y2, color=CYAN, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Editable layered background: large calm technology field.
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY, None)
    add_rect(slide, 0, 0, 13.33, 0.12, CYAN, None)
    add_shape(slide, MSO_SHAPE.ARC, 7.5, 0.82, 4.8, 4.8, NAVY, RGBColor(33, 88, 125), 1.4)
    add_shape(slide, MSO_SHAPE.ARC, 8.2, 1.45, 3.4, 3.4, NAVY, RGBColor(24, 111, 145), 1.1)
    add_shape(slide, MSO_SHAPE.OVAL, 10.16, 3.03, 0.16, 0.16, CYAN, None)
    add_shape(slide, MSO_SHAPE.OVAL, 8.44, 2.08, 0.12, 0.12, BLUE, None)
    add_shape(slide, MSO_SHAPE.OVAL, 11.32, 1.94, 0.12, 0.12, CYAN, None)
    for x1, y1, x2, y2 in [(8.5, 2.14, 10.24, 3.1), (10.24, 3.1, 11.38, 2.0), (8.5, 2.14, 11.38, 2.0)]:
        add_line(slide, x1, y1, x2, y2, RGBColor(40, 117, 152), 0.8)

    add_round(slide, 0.64, 0.72, 1.9, 0.38, RGBColor(18, 61, 96), None)
    add_text(slide, 0.78, 0.815, 1.6, 0.15, "工程实践结题汇报", 8.8, CYAN, True, PP_ALIGN.CENTER)
    add_text(slide, 0.62, 1.52, 6.9, 1.1, "AIChemMCP\n智能调度原型系统", 34, WHITE, True)
    add_text(slide, 0.66, 2.92, 5.65, 0.32,
             "面向自动化化学实验流程的仿真调度与模拟设备接口原型", 12.5, RGBColor(210, 228, 240))
    add_chip(slide, 0.66, 3.62, 1.35, "AIMCP", BLUE)
    add_chip(slide, 2.15, 3.62, 1.68, "Action Server", RGBColor(24, 87, 135))
    add_chip(slide, 3.98, 3.62, 1.88, "Scheduling Runtime", CYAN, NAVY)

    # Hero system panel, fewer words and stronger center.
    add_round(slide, 7.34, 4.62, 4.9, 1.35, RGBColor(14, 50, 82), RGBColor(46, 133, 165), 1.0)
    add_text(slide, 7.62, 4.92, 4.35, 0.28, "Agent -> Runtime -> Scheduler", 17, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 7.70, 5.34, 4.2, 0.22, "任务提交 · 资源调度 · 状态反馈 · 结果归档", 9.8, RGBColor(196, 225, 238), False, PP_ALIGN.CENTER)
    add_text(slide, 7.58, 6.12, 4.48, 0.25, "当前定位：仿真调度原型，未接入真实实验室硬件", 8.6, RGBColor(192, 214, 227), False, PP_ALIGN.CENTER)
    add_footer(slide, dark=True)


def make_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_rect(slide, 0, 0, 13.33, 0.12, CYAN, None)
    add_title(slide, "系统总体架构", "接口层与调度层解耦，形成可运行的仿真调度闭环")

    # Soft central stage.
    add_round(slide, 0.52, 1.35, 12.28, 2.45, RGBColor(238, 246, 252), RGBColor(218, 232, 242), 0.8)
    add_text(slide, 0.86, 1.58, 2.0, 0.18, "主调用链", 8.5, BLUE, True)

    y = 2.02
    w = 1.46
    h = 0.92
    gap = 0.24
    start = 0.76
    chain = [
        ("用户 / Agent", "发起调用", "agent.py"),
        ("Action Server", "工具路由", "servers/action_server.py"),
        ("Tools 封装", "参数转换", "tools/action_server_tools.py"),
        ("Runtime", "维护状态", "scheduling/runtime.py"),
        ("Scheduler", "资源预留", "scheduling/scheduler.py"),
        ("Models", "任务资源", "scheduling/models.py"),
        ("demo_result", "结果归档", "outputs/demo_result.json"),
    ]
    for idx, item in enumerate(chain):
        x = start + idx * (w + gap)
        accent = BLUE if idx in (1, 3, 4) else CYAN
        add_main_card(slide, x, y, w, h, *item, accent=accent)
        if idx < len(chain) - 1:
            add_arrow(slide, x + w + 0.04, y + 0.36, gap - 0.02, 0.18, CYAN)

    # Three smaller support cards.
    add_aux_card(slide, 0.92, 4.25, 3.25, 0.98, "输入数据", "examples/*.json", "示例任务、资源、工作流", BLUE)
    add_aux_card(slide, 5.04, 4.25, 3.25, 0.98, "调度能力", "任务提交 / 状态查询", "时间推进、执行记录", CYAN)
    add_aux_card(slide, 9.16, 4.25, 3.25, 0.98, "验证闭环", "tests + demo smoke", "输出结果可复现", GREEN)

    # Light support lines that do not dominate.
    add_line(slide, 2.55, 4.22, 2.55, 3.77, RGBColor(172, 202, 222), 0.8)
    add_line(slide, 6.65, 4.22, 6.65, 3.77, RGBColor(172, 202, 222), 0.8)
    add_line(slide, 10.75, 4.22, 10.75, 3.77, RGBColor(172, 202, 222), 0.8)

    # Bottom layer band.
    layers = [
        ("接口层", "工具路由"),
        ("运行时层", "状态维护"),
        ("调度层", "资源预留"),
        ("数据模型层", "任务资源"),
        ("验证层", "测试归档"),
    ]
    x = 0.72
    for i, (name, desc) in enumerate(layers):
        fill = NAVY if i == 0 else RGBColor(226, 238, 247)
        title_color = WHITE if i == 0 else NAVY
        desc_color = RGBColor(200, 224, 238) if i == 0 else MID
        add_round(slide, x, 6.08, 2.25, 0.54, fill, None)
        add_text(slide, x + 0.14, 6.18, 0.82, 0.16, name, 8.8, title_color, True)
        add_text(slide, x + 1.05, 6.18, 0.95, 0.16, desc, 8.2, desc_color)
        x += 2.42

    add_text(slide, 0.72, 6.78, 5.8, 0.2, "说明：当前系统为模拟设备接口与仿真调度闭环，不声明真实硬件接入。", 8.5, MID)
    add_footer(slide)


def make_demo_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    add_rect(slide, 0, 0, 13.33, 0.12, CYAN, None)
    add_title(slide, "Demo 演示流程 / 调度闭环", "从示例数据到结构化输出，验证原型可运行")

    # Main demo flow, fewer details.
    add_round(slide, 0.62, 1.45, 12.08, 1.78, WHITE, RGBColor(213, 228, 240), 0.9)
    steps = [
        "加载示例数据", "提交实验任务", "分配资源", "推进模拟时间", "查询状态", "输出结果"
    ]
    xs = [0.95, 2.9, 4.82, 6.67, 8.56, 10.45]
    for i, label in enumerate(steps):
        add_shape(slide, MSO_SHAPE.OVAL, xs[i], 1.83, 0.54, 0.54, BLUE if i in (1, 5) else CYAN, None)
        add_text(slide, xs[i] + 0.13, 1.98, 0.28, 0.12, str(i + 1), 8.5, WHITE if i != 0 else NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, xs[i] - 0.38, 2.53, 1.28, 0.22, label, 9.6, NAVY, True, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(slide, xs[i] + 0.63, 2.01, 1.04, 0.14, CYAN)

    # Result cards.
    add_aux_card(slide, 1.04, 4.05, 3.2, 1.06, "任务状态变化", "4 个任务全部完成", "TASK-0004 为三步 workflow", BLUE)
    add_aux_card(slide, 5.06, 4.05, 3.2, 1.06, "模拟时间推进", "final time = 152", "中途 time 39 产生完成事件", CYAN)
    add_aux_card(slide, 9.08, 4.05, 3.2, 1.06, "输出结果归档", "demo_result.json", "all_completed: true", GREEN)

    add_round(slide, 3.2, 5.72, 6.92, 0.58, RGBColor(235, 245, 250), None)
    add_text(slide, 3.48, 5.88, 6.35, 0.18,
             "真实依据：AIChemMCP-main/outputs/demo_result.json；当前为 simulation / mock hardware。",
             8.8, MID, False, PP_ALIGN.CENTER)
    add_footer(slide)


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = inch(13.333333)
    prs.slide_height = inch(7.5)
    make_cover(prs)
    make_architecture(prs)
    make_demo_loop(prs)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build()
