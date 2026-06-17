from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "AIChemMCP_机器人调度结题汇报_v3.pptx"
SPEECH_PATH = OUT_DIR / "speech.md"


NAVY = RGBColor(9, 26, 48)
NAVY_2 = RGBColor(12, 45, 78)
BLUE = RGBColor(30, 107, 255)
CYAN = RGBColor(18, 196, 218)
DARK = RGBColor(37, 49, 64)
MID = RGBColor(91, 108, 128)
LIGHT_BG = RGBColor(246, 249, 253)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(226, 240, 249)
PALE_CYAN = RGBColor(229, 250, 252)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(245, 146, 64)
RED = RGBColor(222, 88, 88)


def inch(v):
    return Inches(v)


def font(run, size, color=DARK, bold=False, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, x, y, w, h, value, size=14, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    font(run, size, color, bold)
    return box


def shape(slide, shape_type, x, y, w, h, fill, line=None, width=1.0):
    s = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(width)
    return s


def rect(slide, x, y, w, h, fill, line=None):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)


def round_box(slide, x, y, w, h, fill=WHITE, line=PALE, width=1.0):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, width)


def oval(slide, x, y, w, h, fill, line=None):
    return shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill, line)


def line(slide, x1, y1, x2, y2, color=CYAN, width=1.0):
    l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    l.line.color.rgb = color
    l.line.width = Pt(width)
    return l


def arrow(slide, x, y, w, h=0.18, color=CYAN):
    return shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h, color, None)


def top_bar(slide):
    rect(slide, 0, 0, 13.333, 0.12, CYAN, None)


def title(slide, value, subtitle=None, dark=False):
    text(slide, 0.62, 0.38, 6.3, 0.44, value, 25 if len(value) < 14 else 23,
         WHITE if dark else NAVY, True)
    if subtitle:
        text(slide, 0.66, 0.86, 8.6, 0.26, subtitle, 10.2, RGBColor(194, 215, 230) if dark else MID)


def footer(slide, dark=False):
    text(slide, 0.62, 7.12, 6.9, 0.22,
         "AIChemMCP scheduling prototype | simulation / mock hardware",
         7.5, RGBColor(185, 207, 224) if dark else MID)


def chip(slide, x, y, w, label, fill, color=WHITE):
    round_box(slide, x, y, w, 0.32, fill, None)
    text(slide, x + 0.05, y + 0.065, w - 0.1, 0.15, label, 8.2, color, True, PP_ALIGN.CENTER)


def card(slide, x, y, w, h, head, body, fill=WHITE, accent=CYAN, head_size=13):
    round_box(slide, x + 0.035, y + 0.035, w, h, RGBColor(218, 227, 236), None)
    round_box(slide, x, y, w, h, fill, RGBColor(206, 224, 237), 0.9)
    rect(slide, x, y, 0.1, h, accent, None)
    text(slide, x + 0.2, y + 0.16, w - 0.35, 0.24, head, head_size, NAVY, True)
    if body:
        text(slide, x + 0.2, y + 0.48, w - 0.35, h - 0.55, body, 9.0, MID)


def placeholder(slide, x, y, w, h, label):
    round_box(slide, x, y, w, h, RGBColor(239, 245, 250), RGBColor(175, 197, 214), 1.0)
    text(slide, x + 0.18, y + h / 2 - 0.1, w - 0.36, 0.2, label, 9.5, MID, False, PP_ALIGN.CENTER)


def code_panel(slide, x, y, w, h, title_text, lines):
    round_box(slide, x, y, w, h, RGBColor(18, 36, 58), RGBColor(36, 110, 145), 0.9)
    text(slide, x + 0.18, y + 0.14, w - 0.36, 0.22, title_text, 10.2, CYAN, True)
    content = "\n".join(lines)
    box = text(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.62, content, 7.8, RGBColor(220, 235, 244))
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Consolas"
    return box


def intro_dark_background(slide):
    bg(slide, NAVY)
    rect(slide, 0, 0, 13.333, 0.12, CYAN, None)
    shape(slide, MSO_SHAPE.ARC, 8.1, 0.88, 4.2, 4.2, NAVY, RGBColor(34, 91, 126), 1.2)
    shape(slide, MSO_SHAPE.ARC, 8.72, 1.45, 3.0, 3.0, NAVY, RGBColor(23, 118, 150), 1.0)
    for x, y, c in [(9.0, 2.1, BLUE), (10.55, 3.02, CYAN), (11.64, 2.02, CYAN)]:
        oval(slide, x, y, 0.13, 0.13, c, None)
    line(slide, 9.06, 2.16, 10.61, 3.08, RGBColor(40, 117, 152), 0.8)
    line(slide, 10.61, 3.08, 11.70, 2.08, RGBColor(40, 117, 152), 0.8)
    line(slide, 9.06, 2.16, 11.70, 2.08, RGBColor(40, 117, 152), 0.8)


def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    intro_dark_background(s)
    round_box(s, 0.64, 0.72, 1.9, 0.38, RGBColor(18, 61, 96), None)
    text(s, 0.78, 0.815, 1.6, 0.15, "工程实践结题汇报", 8.8, CYAN, True, PP_ALIGN.CENTER)
    text(s, 0.62, 1.52, 6.9, 1.1, "AIChemMCP\n智能调度原型系统", 34, WHITE, True)
    text(s, 0.66, 2.92, 6.35, 0.32, "面向自动化实验室场景的机器人任务调度与模拟设备接口原型", 12.5, RGBColor(210, 228, 240))
    chip(s, 0.66, 3.62, 1.35, "AIChemMCP", BLUE)
    chip(s, 2.15, 3.62, 1.68, "Action Server", RGBColor(24, 87, 135))
    chip(s, 3.98, 3.62, 1.88, "Scheduling Runtime", CYAN, NAVY)
    round_box(s, 7.34, 4.62, 4.9, 1.35, RGBColor(14, 50, 82), RGBColor(46, 133, 165), 1.0)
    text(s, 7.62, 4.92, 4.35, 0.28, "Agent -> Runtime -> Scheduler", 17, WHITE, True, PP_ALIGN.CENTER)
    text(s, 7.70, 5.34, 4.2, 0.22, "任务提交 · 资源调度 · 状态反馈 · 结果归档", 9.8, RGBColor(196, 225, 238), False, PP_ALIGN.CENTER)
    text(s, 7.58, 6.12, 4.48, 0.25, "当前定位：仿真调度原型，未接入真实实验室硬件", 8.6, RGBColor(192, 214, 227), False, PP_ALIGN.CENTER)
    footer(s, True)


def slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "汇报框架", "按照“需求—设计—实现—测试—总结”的工程过程展开")
    items = [
        ("01", "项目背景与目标", "为什么做调度集成\n本项目完成什么"),
        ("02", "系统设计", "架构、调用链\n模型与状态流转"),
        ("03", "实现与演示", "接口、示例数据\nDemo 闭环与输出"),
        ("04", "测试与总结", "测试验证、工程化\n不足与后续扩展"),
    ]
    for i, (num, head, body) in enumerate(items):
        x = 0.85 + (i % 2) * 6.0
        y = 1.65 + (i // 2) * 2.2
        round_box(s, x, y, 5.25, 1.45, WHITE, RGBColor(205, 222, 236), 0.9)
        text(s, x + 0.25, y + 0.22, 0.7, 0.3, num, 18, CYAN if i else BLUE, True)
        text(s, x + 1.05, y + 0.23, 3.4, 0.25, head, 16, NAVY, True)
        text(s, x + 1.05, y + 0.62, 3.7, 0.44, body, 10, MID)
    line(s, 6.66, 1.4, 6.66, 6.25, RGBColor(198, 218, 232), 0.8)
    line(s, 0.92, 3.76, 12.12, 3.76, RGBColor(198, 218, 232), 0.8)
    footer(s)


def slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "项目背景", "自动化实验室任务需要从“计划”走向“可调度执行”")
    # Left concept visual
    round_box(s, 0.75, 1.55, 5.2, 4.6, RGBColor(235, 246, 250), RGBColor(211, 228, 238), 0.8)
    text(s, 1.05, 1.9, 4.4, 0.35, "自动化实验室场景", 18, NAVY, True, PP_ALIGN.CENTER)
    nodes = [("任务目标", 1.25, 3.0), ("任务队列", 3.0, 2.45), ("设备资源", 4.45, 3.45), ("状态反馈", 2.75, 4.55)]
    for label, x, y in nodes:
        oval(s, x, y, 1.05, 0.55, WHITE, CYAN)
        text(s, x + 0.08, y + 0.17, 0.9, 0.13, label, 9.5, NAVY, True, PP_ALIGN.CENTER)
    for a in [(1.95, 3.25, 3.05, 2.75), (3.75, 2.72, 4.5, 3.65), (4.55, 3.92, 3.3, 4.65), (2.8, 4.65, 1.85, 3.52)]:
        line(s, *a, CYAN, 1.1)
    card(s, 6.55, 1.62, 2.7, 1.05, "任务不是单点动作", "示例任务需要排队、调度和状态跟踪", WHITE, BLUE)
    card(s, 9.55, 1.62, 2.7, 1.05, "资源存在竞争", "机器人与工作站需要协同调度", WHITE, CYAN)
    card(s, 6.55, 3.18, 2.7, 1.05, "需要工程闭环", "接口、运行时、输出归档和测试验证", WHITE, CYAN)
    card(s, 9.55, 3.18, 2.7, 1.05, "先做仿真原型", "为后续硬件协议接入预留结构", WHITE, GREEN)
    placeholder(s, 6.55, 5.0, 5.7, 0.75, "此处可补充项目背景或课程要求截图")
    footer(s)


def slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "项目目标与建设内容", "构建可运行、可验证、可扩展的实验室机器人调度仿真原型")
    goals = [
        ("机器人任务抽象", "示例任务步骤\n封装为可调度任务", BLUE),
        ("实验室资源抽象", "Workstation / Robot / Tool\n描述仿真实验室资源", CYAN),
        ("调度运行时", "维护任务、时间、状态\n和执行历史", CYAN),
        ("验证闭环", "状态查询 / 时间推进\ndemo 输出与测试", GREEN),
    ]
    for i, (h, b, c) in enumerate(goals):
        card(s, 0.75 + i * 3.05, 1.75, 2.55, 1.55, h, b, WHITE, c, 15)
    round_box(s, 2.25, 4.35, 8.9, 1.15, RGBColor(235, 246, 250), RGBColor(205, 225, 238), 0.8)
    text(s, 2.65, 4.63, 8.1, 0.32, "目标边界：当前展示的是 simulation / mock hardware 调度闭环", 17, NAVY, True, PP_ALIGN.CENTER)
    text(s, 2.8, 5.07, 7.8, 0.22, "不声明真实机械臂、真实仪器或真实实验执行已经接入", 9.8, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "需求分析", "机器人任务调度、模拟设备接口和演示验证共同约束工程实现")
    cols = [
        ("功能需求", "任务提交\n状态查询\n时间推进\n运行到完成", BLUE),
        ("非功能需求", "结构清晰\n可复现\n错误可返回\n结果可追踪", CYAN),
        ("展示需求", "一键 demo\n结构化输出\n测试结果\n截图占位", GREEN),
    ]
    for i, (h, b, c) in enumerate(cols):
        x = 0.95 + i * 4.05
        round_box(s, x, 1.72, 3.25, 4.3, WHITE, RGBColor(206, 224, 237), 0.9)
        oval(s, x + 0.25, 2.03, 0.46, 0.46, c, None)
        text(s, x + 0.86, 2.07, 1.7, 0.24, h, 16, NAVY, True)
        for j, item in enumerate(b.split("\n")):
            chip(s, x + 0.52, 2.82 + j * 0.58, 2.22, item, RGBColor(236, 246, 252), NAVY)
    text(s, 1.05, 6.45, 10.8, 0.24, "实现原则：先完成可运行的软件原型，再为后续真实硬件接入保留边界。", 9.5, MID, False, PP_ALIGN.CENTER)
    footer(s)


def main_card(slide, x, y, w, h, head, role, path, accent=CYAN):
    round_box(slide, x + 0.035, y + 0.035, w, h, RGBColor(218, 227, 236), None)
    round_box(slide, x, y, w, h, WHITE, RGBColor(191, 214, 230), 1.1)
    rect(slide, x, y, 0.1, h, accent, None)
    text(slide, x + 0.2, y + 0.15, w - 0.35, 0.22, head, 11.5, NAVY, True)
    text(slide, x + 0.2, y + 0.43, w - 0.35, 0.20, role, 8.8, DARK)
    text(slide, x + 0.2, y + 0.68, w - 0.35, 0.18, path, 7.0, MID)


def slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "系统总体架构", "接口层与调度层解耦，形成可运行的仿真调度闭环")
    round_box(s, 0.52, 1.35, 12.28, 2.45, RGBColor(238, 246, 252), RGBColor(218, 232, 242), 0.8)
    text(s, 0.86, 1.58, 2.0, 0.18, "主调用链", 8.5, BLUE, True)
    chain = [
        ("用户 / Agent", "发起调用", "agent.py"),
        ("Action Server", "工具路由", "servers/action_server.py"),
        ("Tools 封装", "参数转换", "tools/action_server_tools.py"),
        ("Runtime", "维护状态", "scheduling/runtime.py"),
        ("Scheduler", "资源预留", "scheduling/scheduler.py"),
        ("Models", "任务资源", "scheduling/models.py"),
        ("demo_result", "结果归档", "outputs/demo_result.json"),
    ]
    start, y, w, gap = 0.76, 2.02, 1.46, 0.24
    for idx, item in enumerate(chain):
        x = start + idx * (w + gap)
        main_card(s, x, y, w, 0.92, *item, accent=BLUE if idx in (1, 3, 4) else CYAN)
        if idx < len(chain) - 1:
            arrow(s, x + w + 0.04, y + 0.36, gap - 0.02, 0.18, CYAN)
    aux = [
        ("输入数据", "examples/*.json", "示例任务、资源、工作流", BLUE),
        ("调度能力", "任务提交 / 状态查询", "时间推进、执行记录", CYAN),
        ("验证闭环", "tests + demo smoke", "输出结果可复现", GREEN),
    ]
    for i, (a, b, c, d) in enumerate(aux):
        x = 0.92 + i * 4.12
        card(s, x, 4.25, 3.25, 0.98, a, b + "\n" + c, WHITE, d, 11)
    layers = [("接口层", "工具路由"), ("运行时层", "状态维护"), ("调度层", "资源预留"), ("数据模型层", "任务资源"), ("验证层", "测试归档")]
    x = 0.72
    for i, (name, desc) in enumerate(layers):
        fill = NAVY if i == 0 else RGBColor(226, 238, 247)
        tc = WHITE if i == 0 else NAVY
        dc = RGBColor(200, 224, 238) if i == 0 else MID
        round_box(s, x, 6.08, 2.25, 0.54, fill, None)
        text(s, x + 0.14, 6.18, 0.82, 0.16, name, 8.8, tc, True)
        text(s, x + 1.05, 6.18, 0.95, 0.16, desc, 8.2, dc)
        x += 2.42
    footer(s)


def slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "核心调用链", "从工具请求到调度状态反馈的闭环")
    steps = [
        ("用户 / Agent", "发起工具调用"),
        ("Action Server", "解析 JSON-RPC"),
        ("Tools", "转为 Runtime 方法"),
        ("Runtime", "提交任务 / 推进时间"),
        ("Scheduler", "生成 START_PROCESSING"),
        ("Status", "返回状态快照"),
    ]
    for i, (h, b) in enumerate(steps):
        x = 1.05 + (i % 3) * 4.0
        y = 1.55 + (i // 3) * 2.1
        card(s, x, y, 2.9, 1.18, h, b, WHITE, BLUE if i < 2 else CYAN, 14)
        if i % 3 < 2:
            arrow(s, x + 3.0, y + 0.48, 0.7, 0.16, CYAN)
    line(s, 10.6, 2.75, 10.6, 3.45, CYAN, 1.2)
    arrow(s, 10.25, 3.35, 0.7, 0.16, CYAN)
    round_box(s, 2.05, 6.1, 9.2, 0.56, RGBColor(234, 245, 250), None)
    text(s, 2.3, 6.26, 8.7, 0.18, "真实文件：servers/action_server.py / tools/action_server_tools.py / scheduling/runtime.py", 8.5, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "调度模型设计", "用统一数据模型承接机器人任务和仿真资源")
    center_x, center_y = 5.2, 2.45
    card(s, center_x, center_y, 2.8, 1.25, "Task", "workflow_tools\nprocessing_times\nstatus / current_step", RGBColor(235, 246, 250), BLUE, 16)
    entities = [
        ("Resource", "timeline\ncurrent_task_id", 1.1, 1.7, CYAN),
        ("Workstation", "tools\nreaction / measure / char", 9.4, 1.7, CYAN),
        ("Robot", "transfer timeline\nRB_1 / RB_2", 1.1, 4.4, GREEN),
        ("Tool", "reaction_tool\nhplc_tool", 9.4, 4.4, GREEN),
    ]
    for h, b, x, y, c in entities:
        card(s, x, y, 2.7, 1.1, h, b, WHITE, c, 14)
        line(s, center_x + 1.4, center_y + 0.62, x + 1.35, y + 0.55, RGBColor(158, 190, 212), 0.9)
    code_panel(s, 4.35, 4.65, 4.45, 1.28, "scheduling/models.py",
               ["TaskStatus: WAITING / RUNNING / COMPLETED / ERROR",
                "ResourceStatus: IDLE / BUSY / RESERVED / TRANSPORTING"])
    footer(s)


def slide_09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "任务状态流转", "运行时维护任务状态与资源状态的可观测变化")
    states = [("WAITING", 1.0, 2.1, BLUE), ("RUNNING", 4.1, 2.1, CYAN), ("COMPLETED", 7.2, 2.1, GREEN), ("ERROR", 10.3, 2.1, ORANGE)]
    for st, x, y, c in states:
        round_box(s, x, y, 2.0, 0.8, WHITE, c, 1.2)
        text(s, x + 0.2, y + 0.28, 1.6, 0.15, st, 13, NAVY, True, PP_ALIGN.CENTER)
    arrow(s, 3.05, 2.43, 0.75, 0.16, CYAN)
    arrow(s, 6.15, 2.43, 0.75, 0.16, CYAN)
    arrow(s, 9.25, 2.43, 0.75, 0.16, ORANGE)
    text(s, 3.05, 2.85, 0.9, 0.15, "调度命令", 7.8, MID, False, PP_ALIGN.CENTER)
    text(s, 6.18, 2.85, 0.85, 0.15, "步骤完成", 7.8, MID, False, PP_ALIGN.CENTER)
    card(s, 1.0, 4.1, 3.15, 1.05, "任务状态", "WAITING -> RUNNING -> COMPLETED\nERROR 为预留状态", WHITE, BLUE, 13)
    card(s, 5.1, 4.1, 3.15, 1.05, "资源状态", "IDLE / BUSY / RESERVED\nMOVING_TO_PICKUP / TRANSPORTING", WHITE, CYAN, 13)
    card(s, 9.2, 4.1, 3.15, 1.05, "状态快照", "get_runtime_status()\n输出任务、资源、历史", WHITE, GREEN, 13)
    footer(s)


def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "工具接口设计", "Action Server 对外暴露 6 个原型系统工具")
    left = [("robotic_reaction", "提交示例反应步骤"),
            ("robotic_measurement", "提交示例测量步骤"),
            ("robotic_characterization", "提交示例表征步骤")]
    right = [("scheduler_status", "查询状态"),
             ("scheduler_advance", "推进时间"),
             ("scheduler_run_until_complete", "运行到完成")]
    round_box(s, 0.82, 1.42, 5.6, 4.55, WHITE, RGBColor(205, 222, 236), 0.9)
    round_box(s, 6.92, 1.42, 5.6, 4.55, WHITE, RGBColor(205, 222, 236), 0.9)
    rect(s, 0.82, 1.42, 5.6, 0.62, RGBColor(232, 244, 252), None)
    rect(s, 6.92, 1.42, 5.6, 0.62, RGBColor(229, 250, 252), None)
    oval(s, 1.15, 1.61, 0.26, 0.26, BLUE, None)
    oval(s, 7.25, 1.61, 0.26, 0.26, CYAN, None)
    text(s, 1.55, 1.59, 2.1, 0.22, "任务提交类", 16, NAVY, True)
    text(s, 7.65, 1.59, 2.1, 0.22, "调度控制类", 16, NAVY, True)
    text(s, 3.85, 1.66, 1.7, 0.16, "3 tools", 8.2, MID, False, PP_ALIGN.RIGHT)
    text(s, 9.95, 1.66, 1.7, 0.16, "3 tools", 8.2, MID, False, PP_ALIGN.RIGHT)
    for i, (name, desc) in enumerate(left):
        card(s, 1.18, 2.35 + i * 0.92, 4.75, 0.66, name, desc, RGBColor(241, 248, 252), BLUE if i == 0 else CYAN, 11)
    for i, (name, desc) in enumerate(right):
        card(s, 7.28, 2.35 + i * 0.92, 4.75, 0.66, name, desc, RGBColor(241, 248, 252), GREEN if i == 2 else CYAN, 11)
    round_box(s, 2.3, 6.28, 8.7, 0.42, RGBColor(235, 245, 250), None)
    text(s, 2.55, 6.41, 8.2, 0.12, "实现依据：servers/action_server.py 的 AVAILABLE_TOOLS_ACTION 与 capability 声明", 8.2, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "示例数据与资源设计", "使用模拟工作站、机器人和示例 workflow 支撑演示")
    card(s, 0.85, 1.5, 3.3, 1.15, "模拟工作站", "WS_REACTOR_A\nWS_MEASURE_A\nWS_CHAR_A", WHITE, BLUE, 14)
    card(s, 0.85, 3.05, 3.3, 1.15, "模拟机器人", "RB_1：转运样品\nRB_2：备用资源", WHITE, CYAN, 14)
    code_panel(s, 4.62, 1.48, 3.5, 2.72, "sample_workflow.json",
               ['"workflow_tools": [',
                '  "reaction_tool",',
                '  "yield_measurement_tool",',
                '  "hplc_tool"',
                ']',
                '"seamless_steps": [[0, 1]]'])
    card(s, 8.65, 1.5, 3.55, 1.2, "示例 workflow", "reaction -> yield -> hplc\n用于演示的任务步骤", WHITE, GREEN, 14)
    placeholder(s, 8.65, 3.18, 3.55, 1.05, "此处可补充 examples JSON 截图")
    text(s, 1.05, 6.25, 11.0, 0.24, "这些资源均为 simulation/mock 配置，用于验证调度逻辑，不代表真实硬件接入。", 9.2, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "Demo 演示流程", "从示例数据到结构化输出，验证机器人调度原型可运行")
    round_box(s, 0.62, 1.45, 12.08, 1.78, WHITE, RGBColor(213, 228, 240), 0.9)
    steps = ["加载示例数据", "提交任务请求", "分配资源", "推进模拟时间", "查询状态", "输出结果"]
    xs = [0.95, 2.9, 4.82, 6.67, 8.56, 10.45]
    for i, label in enumerate(steps):
        oval(s, xs[i], 1.83, 0.54, 0.54, BLUE if i in (1, 5) else CYAN, None)
        text(s, xs[i] + 0.13, 1.98, 0.28, 0.12, str(i + 1), 8.5, WHITE if i != 0 else NAVY, True, PP_ALIGN.CENTER)
        text(s, xs[i] - 0.38, 2.53, 1.28, 0.22, label, 9.6, NAVY, True, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(s, xs[i] + 0.63, 2.01, 1.04, 0.14, CYAN)
    card(s, 1.04, 4.05, 3.2, 1.06, "任务状态变化", "4 个任务全部完成\nTASK-0004 为三步 workflow", WHITE, BLUE, 13)
    card(s, 5.06, 4.05, 3.2, 1.06, "模拟时间推进", "final time = 152\n中途 time 39 产生完成事件", WHITE, CYAN, 13)
    card(s, 9.08, 4.05, 3.2, 1.06, "输出结果归档", "demo_result.json\nall_completed: true", WHITE, GREEN, 13)
    round_box(s, 3.2, 5.72, 6.92, 0.58, RGBColor(235, 245, 250), None)
    text(s, 3.48, 5.88, 6.35, 0.18, "真实依据：AIChemMCP-main/outputs/demo_result.json；当前为 simulation / mock hardware。", 8.8, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "Demo 运行结果", "结果看板展示任务完成、资源占用与输出归档")
    kpis = [("4", "submitted tasks", BLUE), ("152", "simulation time", CYAN), ("true", "all_completed", GREEN), ("113", "steps_run", ORANGE)]
    for i, (v, label, c) in enumerate(kpis):
        round_box(s, 0.85 + i * 3.05, 1.45, 2.45, 1.15, WHITE, RGBColor(205, 222, 236), 0.9)
        text(s, 1.05 + i * 3.05, 1.67, 2.0, 0.32, v, 25, c, True, PP_ALIGN.CENTER)
        text(s, 1.05 + i * 3.05, 2.13, 2.0, 0.18, label, 8.8, MID, False, PP_ALIGN.CENTER)
    code_panel(s, 0.95, 3.18, 5.3, 2.22, "demo_result.json 摘要",
               ['"demo_mode": "simulation"',
                '"final": {',
                '  "time": 152,',
                '  "all_completed": true',
                '}',
                '"TASK-0004": "COMPLETED"'])
    card(s, 6.85, 3.18, 2.6, 1.0, "工作站结果", "WS_REACTOR_A\nWS_MEASURE_A\nWS_CHAR_A\n最终均 IDLE", WHITE, CYAN, 13)
    card(s, 9.85, 3.18, 2.6, 1.0, "机器人结果", "RB_1 处理 TASK-0004\nRB_2 未占用", WHITE, GREEN, 13)
    placeholder(s, 6.85, 4.75, 5.6, 0.65, "此处可补充 demo 运行终端截图")
    footer(s)


def slide_14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "关键功能实现", "核心实现围绕运行时封装、调度预留和结构化输出")
    code_panel(s, 0.82, 1.45, 5.45, 4.2, "scheduling/runtime.py",
               ["submit_task(...)",
                "  validate workflow_tools",
                "  create Task",
                "  scheduler.add_task(task)",
                "  commands = tick(steps=1)",
                "",
                "run_until_all_complete(max_steps)",
                "  advance until all COMPLETED"])
    cards = [
        ("任务提交", "把实验请求转为 Task", BLUE),
        ("资源映射", "tool -> workstation", CYAN),
        ("时间推进", "tick / advance_time", CYAN),
        ("结果归档", "demo_result.json", GREEN),
    ]
    for i, (h, b, c) in enumerate(cards):
        card(s, 6.85 + (i % 2) * 2.8, 1.6 + (i // 2) * 1.55, 2.35, 1.0, h, b, WHITE, c, 13)
    placeholder(s, 6.85, 4.95, 5.15, 0.65, "此处可补充 runtime / scheduler 代码截图")
    footer(s)


def slide_15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "异常处理与工程健壮性", "输入校验和结构化错误返回让原型更可演示、可测试")
    stages = [("输入请求", "sample_id / steps / workflow_tools"), ("运行时校验", "_require_text / _positive_int"), ("错误包装", "_safe_call / _error_response"), ("结构化返回", "INVALID_INPUT + runtime_status")]
    for i, (h, b) in enumerate(stages):
        x = 0.85 + i * 3.05
        card(s, x, 2.0, 2.45, 1.05, h, b, WHITE, ORANGE if i >= 2 else CYAN, 12)
        if i < 3:
            arrow(s, x + 2.55, 2.43, 0.38, 0.16, CYAN if i < 2 else ORANGE)
    code_panel(s, 1.05, 4.2, 5.1, 1.45, "tools/action_server_tools.py",
               ['return { "ok": False,',
                '  "error": {',
                '    "code": "INVALID_INPUT"',
                '  },',
                '  "runtime_status": ... }'])
    card(s, 7.0, 4.2, 4.8, 1.45, "边界说明", "已实现基础输入校验与错误返回；真实设备 ACK、超时重试和安全联锁仍属于后续扩展。", WHITE, ORANGE, 13)
    footer(s)


def slide_16(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "测试方案", "用标准库 unittest 覆盖模型、调度器、运行时、工具和 demo")
    tests = [
        ("模型", "Task 创建\nworkflow 必填", "TaskModelTests"),
        ("调度", "资源预留\n命令生成", "SchedulerTests"),
        ("运行时", "提交 / 状态\n完成 / 推进", "RuntimeTests"),
        ("错误", "INVALID_INPUT\n结构化返回", "ActionServerTools"),
        ("Demo", "生成结果\n输出归档", "DemoSmokeTests"),
    ]
    for i, (head, body, source) in enumerate(tests):
        x = 0.78 + i * 2.48
        round_box(s, x, 1.65, 2.08, 3.65, WHITE, RGBColor(205, 222, 236), 0.9)
        oval(s, x + 0.74, 2.02, 0.56, 0.56, GREEN if i in (3, 4) else CYAN, None)
        text(s, x + 0.84, 2.18, 0.34, 0.13, str(i + 1), 8.5, WHITE if i != 0 else NAVY, True, PP_ALIGN.CENTER)
        text(s, x + 0.28, 2.85, 1.5, 0.22, head, 15, NAVY, True, PP_ALIGN.CENTER)
        text(s, x + 0.24, 3.28, 1.6, 0.46, body, 9.2, MID, False, PP_ALIGN.CENTER)
        text(s, x + 0.18, 4.46, 1.72, 0.15, source, 6.9, MID, False, PP_ALIGN.CENTER)
    round_box(s, 1.85, 5.9, 9.65, 0.52, RGBColor(235, 245, 250), None)
    text(s, 2.15, 6.07, 9.0, 0.14, "入口：AIChemMCP-main/run_tests.py；测试文件：tests/test_scheduling_runtime.py", 8.8, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_17(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "测试结果", "当前运行 8 个测试用例全部通过")
    round_box(s, 0.95, 1.5, 3.5, 1.45, RGBColor(231, 250, 240), GREEN, 1.2)
    text(s, 1.2, 1.86, 3.0, 0.35, "8 tests", 26, GREEN, True, PP_ALIGN.CENTER)
    text(s, 1.2, 2.34, 3.0, 0.18, "Ran 8 tests ... OK", 9.5, MID, False, PP_ALIGN.CENTER)
    code_panel(s, 5.0, 1.5, 6.7, 2.2, "python run_tests.py",
               ["test_demo_generates_result_file ... ok",
                "test_action_tools_return_structured_error ... ok",
                "test_runtime_submit_status_and_completion ... ok",
                "test_scheduler_basic_reservation ... ok",
                "",
                "Ran 8 tests in 0.023s",
                "OK"])
    placeholder(s, 1.1, 4.45, 10.8, 0.85, "此处可补充测试通过终端截图")
    text(s, 1.15, 6.1, 10.8, 0.22, "测试结果来自本轮本地运行，不涉及外部服务。", 9.0, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_18(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "工程化实践", "从概念验证整理为可运行、可测试、可复现的项目结构")
    practices = [
        ("模块化组织", "scheduling\nservers / tools", BLUE),
        ("接口封装", "Action Server\nTools 返回包装", CYAN),
        ("测试验证", "run_tests.py\n8 个 unittest", GREEN),
        ("文档与归档", "README / docs\noutputs JSON", ORANGE),
    ]
    for i, (h, b, c) in enumerate(practices):
        x = 0.95 + i * 3.0
        round_box(s, x, 1.72, 2.45, 2.05, WHITE, RGBColor(205, 222, 236), 0.9)
        oval(s, x + 0.22, 2.05, 0.34, 0.34, c, None)
        text(s, x + 0.72, 2.05, 1.45, 0.22, h, 14.5, NAVY, True)
        text(s, x + 0.32, 2.62, 1.8, 0.46, b, 9.2, MID, False, PP_ALIGN.CENTER)
    placeholder(s, 1.05, 4.55, 4.7, 0.88, "此处可补充项目目录截图")
    card(s, 6.25, 4.55, 5.85, 0.88, "可复现链路", "python demo.py -> outputs/demo_result.json\npython run_tests.py -> OK", RGBColor(235, 246, 250), GREEN, 13)
    footer(s)


def slide_19(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "项目难点与解决方案", "难点集中在抽象、解耦、仿真闭环和可测试性")
    data = [
        ("任务步骤抽象", "统一 Task / workflow_tools"),
        ("Agent 与调度解耦", "Action Server + Tools 网关"),
        ("无硬件条件下演示", "模拟资源 + 时间线"),
        ("多资源冲突", "Scheduler 预留窗口"),
        ("工程验证", "demo + unittest"),
    ]
    for i, (p, sol) in enumerate(data):
        y = 1.48 + i * 0.88
        round_box(s, 0.92, y, 4.25, 0.62, WHITE, RGBColor(210, 225, 238), 0.7)
        rect(s, 0.92, y, 0.1, 0.62, ORANGE if i < 2 else CYAN, None)
        text(s, 1.18, y + 0.21, 3.5, 0.14, p, 9.8, NAVY, True)
        arrow(s, 5.62, y + 0.24, 0.48, 0.12, CYAN)
        round_box(s, 6.45, y, 4.95, 0.62, RGBColor(239, 248, 252), RGBColor(210, 225, 238), 0.7)
        rect(s, 6.45, y, 0.1, 0.62, GREEN if i == 4 else BLUE, None)
        text(s, 6.72, y + 0.21, 4.1, 0.14, sol, 9.8, NAVY, True)
    text(s, 1.05, 6.35, 10.8, 0.22, "解决方案均落在当前真实代码、demo 和测试中；真实硬件接入仍放在后续扩展。", 8.8, MID, False, PP_ALIGN.CENTER)
    footer(s)


def slide_20(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, LIGHT_BG); top_bar(s)
    title(s, "不足与后续扩展", "当前是仿真原型，后续重点是硬件协议、反馈闭环和可视化")
    stages = [
        ("当前原型", "simulation / mock\n已完成软件闭环", BLUE),
        ("硬件协议", "接入真实设备或\n硬件模拟平台", CYAN),
        ("异常恢复", "ACK / 超时 / 重试\n安全联锁", ORANGE),
        ("复杂调度", "依赖约束\n反馈重规划", CYAN),
        ("可视化平台", "任务队列\n资源时间线", GREEN),
    ]
    for i, (h, b, c) in enumerate(stages):
        x = 0.75 + i * 2.45
        card(s, x, 2.0, 2.05, 1.35, h, b, WHITE, c, 12)
        if i < 4:
            arrow(s, x + 2.1, 2.55, 0.28, 0.14, RGBColor(150, 190, 210))
    round_box(s, 1.15, 4.75, 10.9, 0.85, RGBColor(255, 246, 236), ORANGE, 0.9)
    text(s, 1.45, 4.98, 10.3, 0.22, "不能写成已完成：真实机械臂控制、真实实验室执行、生产级调度、复杂反馈闭环。", 10, DARK, True, PP_ALIGN.CENTER)
    footer(s)


def slide_21(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    intro_dark_background(s)
    text(s, 0.72, 0.72, 2.0, 0.22, "总结", 13, CYAN, True)
    text(s, 0.72, 1.42, 6.8, 0.72, "完成从 AIChemMCP 工具调用\n到调度仿真执行的工程闭环", 27, WHITE, True)
    cards = [
        ("可运行", "demo.py 输出可归档结果"),
        ("可观测", "任务、资源、历史状态可查询"),
        ("可验证", "8 个 unittest 本地通过"),
        ("可扩展", "为真实设备协议预留结构"),
    ]
    for i, (h, b) in enumerate(cards):
        x = 0.8 + i * 3.0
        round_box(s, x, 4.25, 2.35, 1.1, RGBColor(14, 50, 82), RGBColor(46, 133, 165), 0.8)
        text(s, x + 0.18, 4.52, 1.98, 0.22, h, 15, CYAN if i != 0 else WHITE, True, PP_ALIGN.CENTER)
        text(s, x + 0.18, 4.9, 1.98, 0.18, b, 8.5, RGBColor(198, 222, 237), False, PP_ALIGN.CENTER)
    text(s, 3.05, 6.15, 7.2, 0.28, "当前定位始终是仿真调度原型，不夸大为真实硬件系统。", 10, RGBColor(198, 222, 237), False, PP_ALIGN.CENTER)
    footer(s, True)


SPEECH = [
("封面", "各位老师好，我们汇报的项目是 AIChemMCP 智能调度原型系统。这个项目面向自动化实验室场景，重点是机器人任务调度与模拟设备接口；当前定位非常明确：它是仿真调度原型，不是已经接入真实实验室硬件的生产系统。本次汇报重点展示我们如何把 Action Server、调度运行时、示例数据、demo 输出和测试验证串成一个可运行的工程闭环。"),
("汇报框架", "本次汇报按照软件工程过程展开。我们先说明项目背景和目标，再进入需求分析与系统设计，随后展示接口、示例数据和 demo 闭环，最后用测试结果、工程化实践、难点和后续扩展做收束。这样安排的目的，是让大家看到项目不是单点算法展示，而是从需求到验证的完整原型实现。"),
("项目背景", "在自动化实验室场景中，机器人、工作站和仪器资源需要协同完成任务。系统不仅要知道任务目标，还需要管理任务队列、资源占用、状态反馈和结果归档。本项目选择先完成一个可验证的仿真闭环：把示例任务、模拟设备资源和调度运行时连接起来，为后续真实硬件协议接入打基础。"),
("项目目标与建设内容", "项目目标可以概括为四件事：第一，把示例 workflow 中的任务步骤抽象成可调度任务；第二，用工作站、机器人和工具描述仿真实验室资源；第三，构建能长期维护状态的 SchedulingRuntime；第四，通过 Action Server 暴露工具接口并生成可归档 demo 输出。这里需要强调，当前实现的是仿真原型，不声明真实设备已经接入。"),
("需求分析", "需求上我们分成三类。功能需求包括任务提交、状态查询、时间推进和运行到完成；非功能需求强调结构清晰、可复现和错误可返回；展示需求要求能一键 demo、能输出结构化结果、能用测试证明运行链路。这些需求共同约束了我们的目录结构和模块划分。"),
("系统总体架构", "这一页是系统总体架构。主链路从用户或 Agent 发起调用开始，经 Action Server 路由，再由 Tools 层转成运行时方法，进入 SchedulingRuntime 和 Scheduler，最后落到 Task / Resource 模型以及 demo_result.json。周围三个辅助区分别是输入数据、调度能力和验证闭环。这个架构的核心是接口层和调度层解耦。"),
("核心调用链", "核心调用链体现了系统如何跑起来。用户或 Agent 发起工具调用，Action Server 解析请求并找到对应工具，Tools 层进行参数转换和异常包装，Runtime 维护任务与时间状态，Scheduler 生成 START_PROCESSING 命令，最后状态快照返回给上层或写入 demo 结果。"),
("调度模型设计", "模型设计是调度系统的基础。Task 保存 workflow_tools、processing_times、状态和当前步骤；Resource 保存时间线和当前任务；Workstation 和 Robot 都是资源的具体类型；Tool 用来连接任务步骤和工作站能力。这样的设计让任务和资源都能被统一调度和观测。"),
("任务状态流转", "任务状态主要从 WAITING 进入 RUNNING，最后到 COMPLETED；ERROR 是预留的错误状态。资源侧也有 IDLE、BUSY、RESERVED、MOVING_TO_PICKUP、TRANSPORTING 等状态。运行时通过状态快照把这些变化暴露出来，因此 demo 中可以看到任务完成事件和资源占用时间线。"),
("工具接口设计", "Action Server 暴露了六个工具。左侧三个用于提交示例任务步骤，名称中保留 reaction、measurement 和 characterization，是为了对应示例 workflow；右侧三个用于调度控制，分别是查询状态、推进时间和运行到完成。这些工具名和能力声明来自当前的 action_server.py，不是额外虚构的接口。"),
("示例数据与资源设计", "示例数据为 demo 提供了可复现输入。sample_resources.json 中定义了三个模拟工作站和两个模拟机器人；sample_workflow.json 中保留 reaction、yield measurement、HPLC 三步 workflow，它们是用于演示的实验任务类型，并包含 seamless_steps。这里的资源都属于 mock 配置，用于验证调度逻辑。"),
("Demo 演示流程", "Demo 的流程是先加载示例数据，再提交任务请求，由调度器分配资源，然后推进模拟时间、查询状态，最终输出 demo_result.json。这个流程证明当前系统可以从输入数据跑到结构化输出，形成一个软件层面的可运行闭环。"),
("Demo 运行结果", "运行结果来自当前项目的 demo_result.json。结果中可以看到 demo_mode 是 simulation，一共提交了四个任务，最终 simulation_time 为 152，all_completed 为 true。TASK-0004 是三步 workflow，最终也完成。下方预留了终端截图位置，后续汇报时可以补上真实运行画面。"),
("关键功能实现", "关键实现集中在运行时和调度器。Runtime 负责校验输入、创建 Task、提交到 Scheduler，并维护状态快照和执行历史；Scheduler 负责资源时间线预留和命令生成；输出侧则把 demo 过程和最终状态归档到 JSON。这样形成了比较完整的工程闭环。"),
("异常处理与工程健壮性", "为了让原型不只是在正常输入下能跑，我们还做了基础健壮性处理。Runtime 会校验空字符串、非正整数和未知工具；Tools 层用 safe_call 把异常包装成结构化错误，包含 INVALID_INPUT 和 runtime_status。不过真实设备 ACK、超时重试和安全联锁还没有实现，属于后续扩展。"),
("测试方案", "测试方案使用 Python 标准库 unittest，入口是 run_tests.py。测试覆盖模型创建、调度器资源预留、运行时提交与完成、时间推进、非法输入处理、工具层结构化错误，以及 demo smoke test。这个测试组合保证原型的核心链路可以被重复验证。"),
("测试结果", "本轮本地运行 python run_tests.py，结果是 8 个测试全部通过，最终输出 OK。测试结果说明当前代码至少在模型、调度、运行时、错误返回和 demo 输出这些关键点上具备基本可靠性。后续可以在这一页补充终端截图。"),
("工程化实践", "工程化方面，项目已经从概念材料整理成了有入口、有模块、有测试、有输出归档的结构。scheduling 负责调度模型和运行时，servers 与 tools 负责接口封装，examples 提供输入，outputs 保存结果，tests 负责验证。这些都是结题汇报中比较重要的工程价值。"),
("项目难点与解决方案", "项目难点主要有五类：任务步骤如何抽象、Agent 和调度器如何解耦、没有真实硬件时如何展示闭环、多资源如何避免冲突，以及如何证明不是只停留在文档层。对应解决方案分别落在 Task 模型、Action Server 网关、模拟资源时间线、Scheduler 预留机制和 demo 加测试验证中。"),
("不足与后续扩展", "当前不足需要主动说明：系统还没有接入真实设备协议，没有真实 ACK 和异常恢复机制，也没有复杂实验结果反馈驱动的重规划。后续可以沿着硬件协议、异常恢复、复杂依赖调度、可视化平台和 LLM 规划反馈几个方向继续扩展。"),
("总结", "总结一下，本项目完成的是从 AIChemMCP 工具调用到调度仿真执行的工程化闭环。它可以运行、可以观测、可以测试，也为后续真实设备接入预留了结构基础。我们不会把它夸大为真实硬件系统，而是把它定位为一个清晰、可复现的智能调度原型。"),
]


def write_speech():
    parts = ["# AIChemMCP 工程实践结题汇报演讲稿\n"]
    for idx, (name, body) in enumerate(SPEECH, start=1):
        parts.append(f"## Slide {idx}: {name}\n\n{body}\n\n---\n\n注意点：\n- 重点：说明本页在工程过程中的作用。\n- 画面引导：先看主标题，再看核心图形和结果卡片。\n- 边界：涉及硬件时主动说明当前为 simulation / mock hardware。\n")
    SPEECH_PATH.write_text("\n".join(parts), encoding="utf-8")


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = inch(13.333333)
    prs.slide_height = inch(7.5)
    for maker in [
        slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07,
        slide_08, slide_09, slide_10, slide_11, slide_12, slide_13, slide_14,
        slide_15, slide_16, slide_17, slide_18, slide_19, slide_20, slide_21,
    ]:
        maker(prs)
    prs.save(PPTX_PATH)
    write_speech()
    print(PPTX_PATH)
    print(SPEECH_PATH)


if __name__ == "__main__":
    build()
